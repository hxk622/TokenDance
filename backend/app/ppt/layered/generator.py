"""
Layered Slide Generator - 分层幻灯片生成器

生成带有程序化背景的高视觉质量幻灯片。
背景为图像层，文字保持可编辑。
"""

import io
import tempfile
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .backgrounds import BackgroundGenerator, BackgroundStyle
from .decorations import DecorationGenerator, DecorationStyle
from .compositor import LayerCompositor, CompositeSpec, CompositePresets, LayerSpec, TextZone


class LayeredSlideStyle(str, Enum):
    """分层幻灯片样式"""
    
    HERO_TITLE = "hero_title"
    SECTION_HEADER = "section_header"
    VISUAL_IMPACT = "visual_impact"
    MINIMAL_CLEAN = "minimal_clean"
    TECH_MODERN = "tech_modern"
    CUSTOM = "custom"


@dataclass
class LayeredSlideContent:
    """分层幻灯片内容"""
    
    style: LayeredSlideStyle = LayeredSlideStyle.HERO_TITLE
    title: str = ""
    subtitle: str = ""
    body: str = ""
    accent_color: str = "#4a90e2"
    base_color: str = "#1a1a2e"
    custom_spec: Optional[CompositeSpec] = None
    
    # 文字样式
    title_font_size: int = 54
    subtitle_font_size: int = 28
    body_font_size: int = 20
    title_color: str = "#ffffff"
    subtitle_color: str = "#cccccc"
    body_color: str = "#e0e0e0"


class LayeredSlideGenerator:
    """分层幻灯片生成器
    
    生成带有程序化背景图像的幻灯片。
    背景作为图像层嵌入，文字保持可编辑。
    
    使用示例：
        generator = LayeredSlideGenerator()
        
        content = LayeredSlideContent(
            style=LayeredSlideStyle.HERO_TITLE,
            title="Welcome to TokenDance",
            subtitle="AI-Powered Presentation",
            accent_color="#4a90e2"
        )
        
        pptx_path = generator.generate_slide(content, "output.pptx")
    """
    
    # PPT 标准尺寸（16:9）
    SLIDE_WIDTH = 1920
    SLIDE_HEIGHT = 1080
    
    def __init__(self):
        """初始化生成器"""
        self.compositor = LayerCompositor(self.SLIDE_WIDTH, self.SLIDE_HEIGHT)
    
    def generate_slide(
        self,
        content: LayeredSlideContent,
        output_path: Optional[Union[str, Path]] = None
    ) -> str:
        """生成单个分层幻灯片
        
        Args:
            content: 幻灯片内容
            output_path: 输出路径（可选）
            
        Returns:
            str: 生成的 PPTX 文件路径
        """
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Emu(self.SLIDE_WIDTH * 914400 // 96)  # 像素转 EMU
        prs.slide_height = Emu(self.SLIDE_HEIGHT * 914400 // 96)
        
        # 添加空白布局
        blank_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_layout)
        
        # 生成并设置背景
        self._set_slide_background(slide, content)
        
        # 添加文字
        self._add_text_content(slide, content)
        
        # 保存
        if output_path is None:
            output_path = tempfile.mktemp(suffix=".pptx")
        
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        
        return str(output_path)
    
    def generate_slides(
        self,
        contents: List[LayeredSlideContent],
        output_path: Optional[Union[str, Path]] = None
    ) -> str:
        """生成多个分层幻灯片
        
        Args:
            contents: 幻灯片内容列表
            output_path: 输出路径（可选）
            
        Returns:
            str: 生成的 PPTX 文件路径
        """
        prs = Presentation()
        prs.slide_width = Emu(self.SLIDE_WIDTH * 914400 // 96)
        prs.slide_height = Emu(self.SLIDE_HEIGHT * 914400 // 96)
        
        blank_layout = prs.slide_layouts[6]
        
        for content in contents:
            slide = prs.slides.add_slide(blank_layout)
            self._set_slide_background(slide, content)
            self._add_text_content(slide, content)
        
        if output_path is None:
            output_path = tempfile.mktemp(suffix=".pptx")
        
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        
        return str(output_path)
    
    def _get_composite_spec(self, content: LayeredSlideContent) -> CompositeSpec:
        """获取合成规格"""
        if content.custom_spec:
            return content.custom_spec
        
        if content.style == LayeredSlideStyle.HERO_TITLE:
            spec = CompositePresets.hero_title(content.accent_color)
        elif content.style == LayeredSlideStyle.SECTION_HEADER:
            spec = CompositePresets.section_header(content.accent_color)
        elif content.style == LayeredSlideStyle.VISUAL_IMPACT:
            spec = CompositePresets.visual_impact(content.accent_color)
        elif content.style == LayeredSlideStyle.MINIMAL_CLEAN:
            spec = CompositePresets.minimal_clean(content.base_color)
        elif content.style == LayeredSlideStyle.TECH_MODERN:
            spec = CompositePresets.tech_modern(content.accent_color)
        else:
            spec = CompositePresets.hero_title(content.accent_color)
        
        # 覆盖颜色
        spec.base_color = content.base_color
        if content.style != LayeredSlideStyle.MINIMAL_CLEAN:
            spec.accent_color = content.accent_color
        
        return spec
    
    def _set_slide_background(self, slide, content: LayeredSlideContent):
        """设置幻灯片背景"""
        spec = self._get_composite_spec(content)
        
        # 生成合成背景图像
        bg_image = self.compositor.composite(spec)
        
        # 保存为临时文件
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            bg_image.save(f.name, "PNG")
            bg_path = f.name
        
        # 设置为幻灯片背景
        slide.shapes.add_picture(
            bg_path,
            Inches(0),
            Inches(0),
            width=Emu(self.SLIDE_WIDTH * 914400 // 96),
            height=Emu(self.SLIDE_HEIGHT * 914400 // 96)
        )
        
        # 将背景图片移到最底层
        shape = slide.shapes[-1]
        spTree = slide.shapes._spTree
        sp = shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)  # 插入到最前面（底层）
        
        # 清理临时文件
        Path(bg_path).unlink(missing_ok=True)
    
    def _add_text_content(self, slide, content: LayeredSlideContent):
        """添加文字内容"""
        spec = self._get_composite_spec(content)
        text_zones = self.compositor.get_text_safe_area(spec.text_zones)
        
        # 添加标题
        if content.title and "title" in text_zones:
            zone = text_zones["title"]
            self._add_text_box(
                slide,
                content.title,
                zone,
                font_size=content.title_font_size,
                font_color=content.title_color,
                bold=True,
                align=PP_ALIGN.CENTER
            )
        
        # 添加副标题
        if content.subtitle and "subtitle" in text_zones:
            zone = text_zones["subtitle"]
            self._add_text_box(
                slide,
                content.subtitle,
                zone,
                font_size=content.subtitle_font_size,
                font_color=content.subtitle_color,
                align=PP_ALIGN.CENTER
            )
        
        # 添加正文
        if content.body and "body" in text_zones:
            zone = text_zones["body"]
            self._add_text_box(
                slide,
                content.body,
                zone,
                font_size=content.body_font_size,
                font_color=content.body_color,
                align=PP_ALIGN.LEFT
            )
        
        # 处理 main_text 区域（用于某些样式）
        if content.title and "main_text" in text_zones:
            zone = text_zones["main_text"]
            self._add_text_box(
                slide,
                content.title,
                zone,
                font_size=content.title_font_size,
                font_color=content.title_color,
                bold=True,
                align=PP_ALIGN.CENTER
            )
        
        # 处理 section 相关区域
        if content.title and "section_title" in text_zones:
            zone = text_zones["section_title"]
            self._add_text_box(
                slide,
                content.title,
                zone,
                font_size=content.title_font_size,
                font_color=content.title_color,
                bold=True,
                align=PP_ALIGN.LEFT
            )
        
        # 处理 content 区域（minimal_clean 样式）
        if "content" in text_zones:
            zone = text_zones["content"]
            text = content.title
            if content.subtitle:
                text += f"\n\n{content.subtitle}"
            if content.body:
                text += f"\n\n{content.body}"
            
            if text.strip():
                self._add_text_box(
                    slide,
                    text,
                    zone,
                    font_size=content.body_font_size,
                    font_color=content.body_color,
                    align=PP_ALIGN.LEFT
                )
    
    def _add_text_box(
        self,
        slide,
        text: str,
        zone: Tuple[int, int, int, int],
        font_size: int = 24,
        font_color: str = "#ffffff",
        bold: bool = False,
        align = PP_ALIGN.LEFT
    ):
        """添加文本框
        
        Args:
            slide: 幻灯片对象
            text: 文本内容
            zone: (x, y, width, height) 像素坐标
            font_size: 字体大小
            font_color: 字体颜色
            bold: 是否加粗
            align: 对齐方式
        """
        x, y, w, h = zone
        
        # 像素转 EMU
        left = Emu(x * 914400 // 96)
        top = Emu(y * 914400 // 96)
        width = Emu(w * 914400 // 96)
        height = Emu(h * 914400 // 96)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        # 垂直居中
        tf.anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        
        # 设置颜色
        hex_color = font_color.lstrip('#')
        r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        run.font.color.rgb = RGBColor(r, g, b)
    
    def generate_background_image(
        self,
        style: LayeredSlideStyle,
        accent_color: str = "#4a90e2",
        base_color: str = "#1a1a2e",
        output_path: Optional[Union[str, Path]] = None
    ) -> Union[str, Image.Image]:
        """单独生成背景图像（不生成 PPTX）
        
        Args:
            style: 幻灯片样式
            accent_color: 强调色
            base_color: 基础色
            output_path: 输出路径（可选，不提供则返回 Image 对象）
            
        Returns:
            str 或 Image.Image: 文件路径或图像对象
        """
        content = LayeredSlideContent(
            style=style,
            accent_color=accent_color,
            base_color=base_color
        )
        
        spec = self._get_composite_spec(content)
        bg_image = self.compositor.composite(spec)
        
        if output_path:
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            bg_image.save(str(output_path))
            return str(output_path)
        
        return bg_image
    
    def preview_all_styles(
        self,
        output_dir: Union[str, Path],
        accent_color: str = "#4a90e2"
    ) -> List[str]:
        """预览所有样式
        
        Args:
            output_dir: 输出目录
            accent_color: 强调色
            
        Returns:
            List[str]: 生成的图像路径列表
        """
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        paths = []
        
        for style in LayeredSlideStyle:
            if style == LayeredSlideStyle.CUSTOM:
                continue
            
            path = output_dir / f"{style.value}.png"
            self.generate_background_image(
                style=style,
                accent_color=accent_color,
                output_path=path
            )
            paths.append(str(path))
        
        return paths
