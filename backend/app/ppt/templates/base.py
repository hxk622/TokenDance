"""
Slide Template Base Class - 幻灯片模板基类

定义模板的抽象接口，所有具体模板必须继承此类。
"""

from abc import ABC, abstractmethod
from typing import TYPE_CHECKING, Any, Dict, List, Optional, Type

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

if TYPE_CHECKING:
    from ..models import SlideContent, BrandKit, TemplateConfig, SlideType


class SlideTemplate(ABC):
    """幻灯片模板抽象基类
    
    所有模板必须实现：
    - slide_type: 支持的幻灯片类型
    - render(): 渲染幻灯片
    
    可选实现：
    - validate(): 验证内容数据
    - get_preview(): 获取预览信息
    """
    
    # 子类必须定义
    slide_type: "SlideType"
    name: str = ""
    description: str = ""
    
    def __init__(self):
        if not self.name:
            self.name = self.__class__.__name__
    
    @abstractmethod
    def render(
        self,
        prs: Presentation,
        content: "SlideContent",
        brand: "BrandKit",
        config: "TemplateConfig"
    ) -> Slide:
        """渲染幻灯片
        
        Args:
            prs: Presentation 对象
            content: 幻灯片内容
            brand: 品牌配置
            config: 模板配置
            
        Returns:
            Slide: 渲染后的幻灯片
        """
        pass
    
    def validate(self, content: "SlideContent") -> List[str]:
        """验证内容数据
        
        Args:
            content: 幻灯片内容
            
        Returns:
            List[str]: 验证错误列表（空表示通过）
        """
        errors = []
        return errors
    
    def get_preview(self, content: "SlideContent") -> Dict[str, Any]:
        """获取预览信息
        
        Args:
            content: 幻灯片内容
            
        Returns:
            Dict: 预览信息（用于 UI 展示）
        """
        return {
            "type": self.slide_type,
            "template": self.name,
            "title": content.title,
        }
    
    # =========================================================================
    # 辅助方法
    # =========================================================================
    
    @staticmethod
    def hex_to_rgb(hex_color: str) -> RGBColor:
        """将 Hex 颜色转换为 RGBColor
        
        Args:
            hex_color: Hex 颜色代码（如 '#1a1a2e'）
            
        Returns:
            RGBColor: RGB 颜色对象
        """
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    
    @staticmethod
    def add_text_frame(
        shape,
        text: str,
        font_name: str = "Microsoft YaHei",
        font_size: int = 18,
        font_color: str = "#1f2937",
        bold: bool = False,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        vertical: MSO_ANCHOR = MSO_ANCHOR.TOP
    ):
        """为形状添加文本框
        
        Args:
            shape: 目标形状
            text: 文本内容
            font_name: 字体名称
            font_size: 字号
            font_color: 字体颜色
            bold: 是否加粗
            alignment: 水平对齐
            vertical: 垂直对齐
        """
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        # 设置垂直对齐
        tf.vertical_anchor = vertical
        
        # 清除默认段落并添加新内容
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = SlideTemplate.hex_to_rgb(font_color)
        p.font.bold = bold
        p.alignment = alignment
    
    @staticmethod
    def add_bullet_list(
        shape,
        bullets: List[str],
        font_name: str = "Microsoft YaHei",
        font_size: int = 18,
        font_color: str = "#1f2937",
        bullet_color: str = "#4a90e2",
        line_spacing: float = 1.5
    ):
        """为形状添加要点列表
        
        Args:
            shape: 目标形状
            bullets: 要点列表
            font_name: 字体名称
            font_size: 字号
            font_color: 字体颜色
            bullet_color: 项目符号颜色
            line_spacing: 行间距
        """
        tf = shape.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = bullet
            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.font.color.rgb = SlideTemplate.hex_to_rgb(font_color)
            p.level = 0
            
            # 设置行间距
            p.line_spacing = line_spacing
            
            # 设置项目符号
            p.bullet = True
    
    def create_blank_slide(self, prs: Presentation) -> Slide:
        """创建空白幻灯片
        
        Args:
            prs: Presentation 对象
            
        Returns:
            Slide: 空白幻灯片
        """
        # 使用空白布局（索引 6 通常是空白）
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)
    
    def set_background_color(
        self,
        slide: Slide,
        color: str
    ):
        """设置幻灯片背景色
        
        Args:
            slide: 幻灯片
            color: Hex 颜色
        """
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.hex_to_rgb(color)
    
    def add_shape_with_text(
        self,
        slide: Slide,
        left: int,
        top: int,
        width: int,
        height: int,
        text: str,
        **text_kwargs
    ):
        """添加带文本的矩形形状
        
        Args:
            slide: 幻灯片
            left, top, width, height: 位置和尺寸 (EMU)
            text: 文本内容
            **text_kwargs: 传递给 add_text_frame 的参数
            
        Returns:
            Shape: 创建的形状
        """
        from pptx.enum.shapes import MSO_SHAPE
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        # 设置透明背景
        shape.fill.background()
        shape.line.fill.background()
        
        # 添加文本
        self.add_text_frame(shape, text, **text_kwargs)
        
        return shape
    
    def __repr__(self) -> str:
        return f"<{self.__class__.__name__}(type={self.slide_type})>"
