"""
Built-in Slide Templates - 内置幻灯片模板

包含：
- TitleSlideTemplate: 标题页
- SectionSlideTemplate: 章节分隔页
- BulletSlideTemplate: 要点列表页
- TwoColumnSlideTemplate: 双栏布局
- ImageSlideTemplate: 图片页
- ImageTextSlideTemplate: 图文混排
- QuoteSlideTemplate: 引用页
- ContentSlideTemplate: 通用内容页
- ThankYouSlideTemplate: 结束感谢页
- BlankSlideTemplate: 空白页
"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt

from ...models import BrandKit, SlideContent, SlideType, TemplateConfig
from ..base import SlideTemplate

# =============================================================================
# 标题页模板
# =============================================================================

class TitleSlideTemplate(SlideTemplate):
    """标题页模板

    用于演示文稿的第一页，包含主标题和副标题。
    """

    slide_type = SlideType.TITLE
    name = "title"
    description = "标题页，包含主标题和副标题"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.background_color)

        # 计算内容区域
        content_width = config.width - config.margin_left - config.margin_right

        # 主标题 - 居中偏上
        title_top = int(config.height * 0.35)
        title_height = Inches(1.5)

        if content.title:
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                title_top,
                content_width,
                title_height
            )
            title_shape.fill.background()
            title_shape.line.fill.background()

            self.add_text_frame(
                title_shape,
                content.title,
                font_name=brand.title_font,
                font_size=config.title_font_size,
                font_color=brand.primary_color,
                bold=True,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

        # 副标题 - 主标题下方
        if content.subtitle:
            subtitle_top = title_top + title_height + Inches(0.3)
            subtitle_height = Inches(0.8)

            subtitle_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                subtitle_top,
                content_width,
                subtitle_height
            )
            subtitle_shape.fill.background()
            subtitle_shape.line.fill.background()

            self.add_text_frame(
                subtitle_shape,
                content.subtitle,
                font_name=brand.body_font,
                font_size=config.subtitle_font_size,
                font_color=brand.text_secondary,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

        # 添加装饰线
        line_top = int(config.height * 0.65)
        line_width = Inches(2)
        line_left = (config.width - line_width) // 2

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            line_left,
            line_top,
            line_width,
            Pt(4)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.hex_to_rgb(brand.secondary_color)
        line.line.fill.background()

        return slide


# =============================================================================
# 章节分隔页模板
# =============================================================================

class SectionSlideTemplate(SlideTemplate):
    """章节分隔页模板

    用于分隔演示文稿的不同章节。
    """

    slide_type = SlideType.SECTION
    name = "section"
    description = "章节分隔页"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)

        # 使用主色作为背景
        self.set_background_color(slide, brand.primary_color)

        content_width = config.width - config.margin_left - config.margin_right

        # 章节标题 - 居中
        if content.title:
            title_top = int(config.height * 0.4)
            title_height = Inches(1.2)

            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                title_top,
                content_width,
                title_height
            )
            title_shape.fill.background()
            title_shape.line.fill.background()

            # 白色文字
            self.add_text_frame(
                title_shape,
                content.title,
                font_name=brand.title_font,
                font_size=40,
                font_color="#ffffff",
                bold=True,
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

        return slide


# =============================================================================
# 要点列表页模板
# =============================================================================

class BulletSlideTemplate(SlideTemplate):
    """要点列表页模板

    标题 + 要点列表的经典布局。
    """

    slide_type = SlideType.BULLET
    name = "bullet"
    description = "要点列表页"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.background_color)

        content_width = config.width - config.margin_left - config.margin_right

        # 标题
        if content.title:
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                config.margin_top,
                content_width,
                Inches(1)
            )
            title_shape.fill.background()
            title_shape.line.fill.background()

            self.add_text_frame(
                title_shape,
                content.title,
                font_name=brand.title_font,
                font_size=32,
                font_color=brand.primary_color,
                bold=True
            )

        # 要点列表
        if content.bullets:
            bullets_top = config.margin_top + Inches(1.3)
            bullets_height = config.height - bullets_top - config.margin_bottom

            bullets_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left + Inches(0.2),
                bullets_top,
                content_width - Inches(0.4),
                bullets_height
            )
            bullets_shape.fill.background()
            bullets_shape.line.fill.background()

            self.add_bullet_list(
                bullets_shape,
                content.bullets,
                font_name=brand.body_font,
                font_size=config.bullet_font_size,
                font_color=brand.text_primary,
                bullet_color=brand.secondary_color,
                line_spacing=config.line_spacing
            )

        return slide


# =============================================================================
# 双栏布局模板
# =============================================================================

class TwoColumnSlideTemplate(SlideTemplate):
    """双栏布局模板

    左右两栏并列的布局，适合对比或并列展示。
    """

    slide_type = SlideType.TWO_COLUMN
    name = "two_column"
    description = "双栏布局"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.background_color)

        content_width = config.width - config.margin_left - config.margin_right
        column_width = (content_width - Inches(0.5)) // 2

        # 标题
        if content.title:
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                config.margin_top,
                content_width,
                Inches(0.9)
            )
            title_shape.fill.background()
            title_shape.line.fill.background()

            self.add_text_frame(
                title_shape,
                content.title,
                font_name=brand.title_font,
                font_size=32,
                font_color=brand.primary_color,
                bold=True
            )

        # 双栏内容
        columns_top = config.margin_top + Inches(1.2)
        columns_height = config.height - columns_top - config.margin_bottom

        if content.columns and len(content.columns) >= 2:
            # 左栏
            left_col = content.columns[0]
            left_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                columns_top,
                column_width,
                columns_height
            )
            left_shape.fill.solid()
            left_shape.fill.fore_color.rgb = self.hex_to_rgb(brand.surface_color)
            left_shape.line.fill.background()

            # 左栏标题
            tf = left_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = left_col.get("title", "")
            p.font.name = brand.title_font
            p.font.size = Pt(20)
            p.font.color.rgb = self.hex_to_rgb(brand.primary_color)
            p.font.bold = True

            # 左栏内容
            if left_col.get("content"):
                p2 = tf.add_paragraph()
                p2.text = "\n" + left_col.get("content", "")
                p2.font.name = brand.body_font
                p2.font.size = Pt(16)
                p2.font.color.rgb = self.hex_to_rgb(brand.text_primary)

            # 右栏
            right_col = content.columns[1]
            right_left = config.margin_left + column_width + Inches(0.5)
            right_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                right_left,
                columns_top,
                column_width,
                columns_height
            )
            right_shape.fill.solid()
            right_shape.fill.fore_color.rgb = self.hex_to_rgb(brand.surface_color)
            right_shape.line.fill.background()

            # 右栏标题和内容
            tf = right_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = right_col.get("title", "")
            p.font.name = brand.title_font
            p.font.size = Pt(20)
            p.font.color.rgb = self.hex_to_rgb(brand.primary_color)
            p.font.bold = True

            if right_col.get("content"):
                p2 = tf.add_paragraph()
                p2.text = "\n" + right_col.get("content", "")
                p2.font.name = brand.body_font
                p2.font.size = Pt(16)
                p2.font.color.rgb = self.hex_to_rgb(brand.text_primary)

        return slide


# =============================================================================
# 图片页模板
# =============================================================================

class ImageSlideTemplate(SlideTemplate):
    """图片页模板

    全屏或大面积图片展示。
    """

    slide_type = SlideType.IMAGE
    name = "image"
    description = "图片页"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.background_color)

        # 添加图片
        if content.image_path and Path(content.image_path).exists():
            # 居中显示，保持比例
            img_width = config.width - config.margin_left - config.margin_right
            img_height = config.height - config.margin_top - config.margin_bottom - Inches(1)

            slide.shapes.add_picture(
                content.image_path,
                config.margin_left,
                config.margin_top,
                width=img_width,
                height=img_height
            )

        # 图片说明
        if content.image_caption:
            caption_top = config.height - config.margin_bottom - Inches(0.6)
            caption_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                caption_top,
                config.width - config.margin_left - config.margin_right,
                Inches(0.5)
            )
            caption_shape.fill.background()
            caption_shape.line.fill.background()

            self.add_text_frame(
                caption_shape,
                content.image_caption,
                font_name=brand.body_font,
                font_size=14,
                font_color=brand.text_secondary,
                alignment=PP_ALIGN.CENTER
            )

        return slide


# =============================================================================
# 图文混排模板
# =============================================================================

class ImageTextSlideTemplate(SlideTemplate):
    """图文混排模板

    图片 + 文字的混合布局。
    """

    slide_type = SlideType.IMAGE_TEXT
    name = "image_text"
    description = "图文混排"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.background_color)

        content_width = config.width - config.margin_left - config.margin_right
        content_height = config.height - config.margin_top - config.margin_bottom

        # 根据图片位置决定布局
        is_image_left = content.image_position == "left"
        half_width = (content_width - Inches(0.5)) // 2

        # 标题（顶部）
        title_height = Inches(0.9)
        if content.title:
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                config.margin_top,
                content_width,
                title_height
            )
            title_shape.fill.background()
            title_shape.line.fill.background()

            self.add_text_frame(
                title_shape,
                content.title,
                font_name=brand.title_font,
                font_size=32,
                font_color=brand.primary_color,
                bold=True
            )

        # 内容区域
        body_top = config.margin_top + title_height + Inches(0.3)
        body_height = content_height - title_height - Inches(0.3)

        # 图片
        if is_image_left:
            img_left = config.margin_left
            text_left = config.margin_left + half_width + Inches(0.5)
        else:
            img_left = config.margin_left + half_width + Inches(0.5)
            text_left = config.margin_left

        if content.image_path and Path(content.image_path).exists():
            slide.shapes.add_picture(
                content.image_path,
                img_left,
                body_top,
                width=half_width,
                height=body_height
            )

        # 文字内容
        text_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            text_left,
            body_top,
            half_width,
            body_height
        )
        text_shape.fill.background()
        text_shape.line.fill.background()

        # 如果有要点列表，显示要点
        if content.bullets:
            self.add_bullet_list(
                text_shape,
                content.bullets,
                font_name=brand.body_font,
                font_size=config.bullet_font_size,
                font_color=brand.text_primary
            )
        elif content.body:
            self.add_text_frame(
                text_shape,
                content.body,
                font_name=brand.body_font,
                font_size=config.body_font_size,
                font_color=brand.text_primary
            )

        return slide


# =============================================================================
# 引用页模板
# =============================================================================

class QuoteSlideTemplate(SlideTemplate):
    """引用页模板

    突出显示引用文字。
    """

    slide_type = SlideType.QUOTE
    name = "quote"
    description = "引用页"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.surface_color)

        content_width = config.width - config.margin_left - config.margin_right

        # 引用符号
        quote_mark_top = int(config.height * 0.25)
        quote_mark = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            config.margin_left,
            quote_mark_top,
            Inches(1),
            Inches(0.8)
        )
        quote_mark.fill.background()
        quote_mark.line.fill.background()

        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(72)
        p.font.color.rgb = self.hex_to_rgb(brand.secondary_color)

        # 引用内容
        if content.quote:
            quote_top = quote_mark_top + Inches(0.5)
            quote_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left + Inches(0.5),
                quote_top,
                content_width - Inches(1),
                Inches(2)
            )
            quote_shape.fill.background()
            quote_shape.line.fill.background()

            self.add_text_frame(
                quote_shape,
                content.quote,
                font_name=brand.body_font,
                font_size=28,
                font_color=brand.text_primary,
                alignment=PP_ALIGN.LEFT
            )

        # 作者
        if content.quote_author:
            author_top = int(config.height * 0.65)
            author_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left + Inches(0.5),
                author_top,
                content_width - Inches(1),
                Inches(0.6)
            )
            author_shape.fill.background()
            author_shape.line.fill.background()

            self.add_text_frame(
                author_shape,
                f"— {content.quote_author}",
                font_name=brand.body_font,
                font_size=18,
                font_color=brand.text_secondary,
                alignment=PP_ALIGN.RIGHT
            )

        return slide


# =============================================================================
# 通用内容页模板
# =============================================================================

class ContentSlideTemplate(SlideTemplate):
    """通用内容页模板

    标题 + 正文的基础布局。
    """

    slide_type = SlideType.CONTENT
    name = "content"
    description = "通用内容页"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.background_color)

        content_width = config.width - config.margin_left - config.margin_right

        # 标题
        if content.title:
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                config.margin_top,
                content_width,
                Inches(1)
            )
            title_shape.fill.background()
            title_shape.line.fill.background()

            self.add_text_frame(
                title_shape,
                content.title,
                font_name=brand.title_font,
                font_size=32,
                font_color=brand.primary_color,
                bold=True
            )

        # 正文
        body_top = config.margin_top + Inches(1.3)
        body_height = config.height - body_top - config.margin_bottom

        if content.body:
            body_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                body_top,
                content_width,
                body_height
            )
            body_shape.fill.background()
            body_shape.line.fill.background()

            self.add_text_frame(
                body_shape,
                content.body,
                font_name=brand.body_font,
                font_size=config.body_font_size,
                font_color=brand.text_primary
            )

        return slide


# =============================================================================
# 结束感谢页模板
# =============================================================================

class ThankYouSlideTemplate(SlideTemplate):
    """结束感谢页模板

    演示文稿的结束页。
    """

    slide_type = SlideType.THANK_YOU
    name = "thank_you"
    description = "结束感谢页"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.primary_color)

        content_width = config.width - config.margin_left - config.margin_right

        # 感谢文字
        title = content.title or "Thank You"
        title_top = int(config.height * 0.4)

        title_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            config.margin_left,
            title_top,
            content_width,
            Inches(1.5)
        )
        title_shape.fill.background()
        title_shape.line.fill.background()

        self.add_text_frame(
            title_shape,
            title,
            font_name=brand.title_font,
            font_size=48,
            font_color="#ffffff",
            bold=True,
            alignment=PP_ALIGN.CENTER,
            vertical=MSO_ANCHOR.MIDDLE
        )

        # 副标题/联系方式
        if content.subtitle:
            subtitle_top = title_top + Inches(1.8)

            subtitle_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                config.margin_left,
                subtitle_top,
                content_width,
                Inches(0.8)
            )
            subtitle_shape.fill.background()
            subtitle_shape.line.fill.background()

            self.add_text_frame(
                subtitle_shape,
                content.subtitle,
                font_name=brand.body_font,
                font_size=20,
                font_color="#cccccc",
                alignment=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE
            )

        return slide


# =============================================================================
# 空白页模板
# =============================================================================

class BlankSlideTemplate(SlideTemplate):
    """空白页模板

    纯净的空白幻灯片。
    """

    slide_type = SlideType.BLANK
    name = "blank"
    description = "空白页"

    def render(
        self,
        prs: Presentation,
        content: SlideContent,
        brand: BrandKit,
        config: TemplateConfig
    ) -> Slide:
        slide = self.create_blank_slide(prs)
        self.set_background_color(slide, brand.background_color)
        return slide


# =============================================================================
# 辅助函数
# =============================================================================

def get_all_templates() -> list[type[SlideTemplate]]:
    """获取所有内置模板类"""
    return [
        TitleSlideTemplate,
        SectionSlideTemplate,
        BulletSlideTemplate,
        TwoColumnSlideTemplate,
        ImageSlideTemplate,
        ImageTextSlideTemplate,
        QuoteSlideTemplate,
        ContentSlideTemplate,
        ThankYouSlideTemplate,
        BlankSlideTemplate,
    ]


__all__ = [
    "TitleSlideTemplate",
    "SectionSlideTemplate",
    "BulletSlideTemplate",
    "TwoColumnSlideTemplate",
    "ImageSlideTemplate",
    "ImageTextSlideTemplate",
    "QuoteSlideTemplate",
    "ContentSlideTemplate",
    "ThankYouSlideTemplate",
    "BlankSlideTemplate",
    "get_all_templates",
]
