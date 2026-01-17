"""
PPT Generator - PPT 生成器

协调模板引擎生成完整的演示文稿并导出。
"""

import logging
from datetime import datetime
from pathlib import Path

from pptx import Presentation

from .models import (
    BrandKit,
    PresentationSpec,
    SlideContent,
    SlideType,
    TemplateConfig,
)
from .template_engine import TemplateEngine, get_global_engine

logger = logging.getLogger(__name__)


class PPTGenerator:
    """PPT 生成器

    负责协调模板引擎生成完整的演示文稿。

    使用示例：
        generator = PPTGenerator()

        # 方式一：使用 PresentationSpec
        spec = PresentationSpec(title="我的演示")
        spec.add_title_slide("欢迎", "副标题")
        spec.add_bullet_slide("要点", ["要点1", "要点2"])
        pptx_path = await generator.generate_from_spec(spec, "output.pptx")

        # 方式二：使用 SlideContent 列表
        slides = [
            SlideContent(type=SlideType.TITLE, title="标题"),
            SlideContent(type=SlideType.BULLET, title="要点", bullets=["1", "2"]),
        ]
        pptx_path = await generator.generate(slides, "output.pptx")
    """

    def __init__(
        self,
        template_engine: TemplateEngine | None = None,
        output_dir: str = "/tmp/tokendance/ppt"
    ):
        """初始化生成器

        Args:
            template_engine: 模板引擎（可选，默认使用全局引擎）
            output_dir: 默认输出目录
        """
        self.engine = template_engine or get_global_engine()
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"PPTGenerator initialized with output_dir: {self.output_dir}")

    async def generate(
        self,
        slides: list[SlideContent],
        output_path: str | None = None,
        brand: BrandKit | None = None,
        config: TemplateConfig | None = None,
        title: str = "Presentation"
    ) -> str:
        """生成 PPT

        Args:
            slides: 幻灯片内容列表
            output_path: 输出路径（可选）
            brand: 品牌配置（可选）
            config: 模板配置（可选）
            title: 演示文稿标题

        Returns:
            str: 生成的 PPTX 文件路径
        """
        brand = brand or BrandKit()
        config = config or TemplateConfig()

        # 创建 Presentation
        prs = Presentation()

        # 设置幻灯片尺寸
        prs.slide_width = config.width
        prs.slide_height = config.height

        logger.info(f"Generating PPT with {len(slides)} slides")

        # 渲染每一页
        for i, slide_content in enumerate(slides):
            try:
                template = self.engine.select_template(slide_content)
                template.render(prs, slide_content, brand, config)
                logger.debug(f"Rendered slide {i+1}: {template.name}")
            except Exception as e:
                logger.error(f"Failed to render slide {i+1}: {e}")
                # 继续渲染其他幻灯片
                continue

        # 设置文档属性
        prs.core_properties.title = title
        prs.core_properties.created = datetime.now()

        # 确定输出路径
        if output_path:
            final_path = Path(output_path)
        else:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
            final_path = self.output_dir / filename

        # 确保目录存在
        final_path.parent.mkdir(parents=True, exist_ok=True)

        # 保存
        prs.save(str(final_path))
        logger.info(f"PPT saved to: {final_path}")

        return str(final_path)

    async def generate_from_spec(
        self,
        spec: PresentationSpec,
        output_path: str | None = None
    ) -> str:
        """从 PresentationSpec 生成 PPT

        Args:
            spec: 演示文稿规格
            output_path: 输出路径（可选）

        Returns:
            str: 生成的 PPTX 文件路径
        """
        return await self.generate(
            slides=spec.slides,
            output_path=output_path,
            brand=spec.brand,
            config=spec.template_config,
            title=spec.title
        )

    async def generate_from_outline(
        self,
        outline: list[dict],
        brand: BrandKit | None = None,
        output_path: str | None = None,
        title: str = "Presentation"
    ) -> str:
        """从大纲生成 PPT

        简化的接口，从结构化大纲直接生成 PPT。

        Args:
            outline: 大纲列表，格式:
                [
                    {"type": "title", "title": "标题", "subtitle": "副标题"},
                    {"type": "bullet", "title": "要点", "bullets": ["1", "2"]},
                    ...
                ]
            brand: 品牌配置
            output_path: 输出路径
            title: 演示文稿标题

        Returns:
            str: 生成的 PPTX 文件路径
        """
        slides = []

        for item in outline:
            slide_type_str = item.get("type", "content")

            # 转换类型字符串
            try:
                slide_type = SlideType(slide_type_str)
            except ValueError:
                slide_type = SlideType.CONTENT

            # 创建 SlideContent
            slide = SlideContent(
                type=slide_type,
                title=item.get("title"),
                subtitle=item.get("subtitle"),
                body=item.get("body"),
                bullets=item.get("bullets"),
                columns=item.get("columns"),
                image_path=item.get("image_path"),
                image_caption=item.get("image_caption"),
                quote=item.get("quote"),
                quote_author=item.get("quote_author"),
                speaker_notes=item.get("speaker_notes"),
            )
            slides.append(slide)

        return await self.generate(
            slides=slides,
            output_path=output_path,
            brand=brand,
            title=title
        )

    def create_quick_presentation(
        self,
        title: str,
        subtitle: str | None = None,
        sections: list[dict] | None = None,
        brand: BrandKit | None = None
    ) -> PresentationSpec:
        """快速创建演示文稿规格

        便捷方法，自动添加标题页和感谢页。

        Args:
            title: 演示标题
            subtitle: 副标题
            sections: 章节列表，格式:
                [
                    {
                        "title": "章节标题",
                        "slides": [
                            {"type": "bullet", "title": "要点", "bullets": [...]},
                        ]
                    }
                ]
            brand: 品牌配置

        Returns:
            PresentationSpec: 演示文稿规格
        """
        spec = PresentationSpec(
            title=title,
            brand=brand or BrandKit()
        )

        # 添加标题页
        spec.add_title_slide(title, subtitle)

        # 添加章节内容
        if sections:
            for section in sections:
                # 章节分隔页
                if section.get("title"):
                    spec.add_section_slide(section["title"])

                # 章节内容页
                for slide_data in section.get("slides", []):
                    slide_type_str = slide_data.get("type", "content")
                    try:
                        slide_type = SlideType(slide_type_str)
                    except ValueError:
                        slide_type = SlideType.CONTENT

                    spec.add_slide(SlideContent(
                        type=slide_type,
                        **{k: v for k, v in slide_data.items() if k != "type"}
                    ))

        # 添加感谢页
        spec.add_slide(SlideContent(
            type=SlideType.THANK_YOU,
            title="Thank You",
            subtitle=subtitle
        ))

        return spec

    def get_supported_slide_types(self) -> list[str]:
        """获取支持的幻灯片类型

        Returns:
            List[str]: 类型列表
        """
        return [t.value for t in SlideType]

    def get_template_info(self) -> list[dict]:
        """获取模板信息

        Returns:
            List[dict]: 模板信息列表
        """
        return self.engine.list_templates()


# 便捷函数
async def quick_generate(
    slides: list[dict],
    output_path: str | None = None,
    title: str = "Presentation"
) -> str:
    """快速生成 PPT

    便捷函数，直接从字典列表生成 PPT。

    Args:
        slides: 幻灯片数据列表
        output_path: 输出路径
        title: 标题

    Returns:
        str: 生成的文件路径
    """
    generator = PPTGenerator()
    return await generator.generate_from_outline(
        outline=slides,
        output_path=output_path,
        title=title
    )
