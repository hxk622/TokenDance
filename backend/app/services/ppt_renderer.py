"""
PPT Renderer Service - PPT 渲染服务

基于 Marp CLI 的 PPT 渲染和导出：
- Markdown → HTML 预览
- Markdown → PDF 导出
- Markdown → PPTX 导出（通过 Puppeteer）
- 主题管理
- 自定义样式支持

依赖：
- @marp-team/marp-cli: npm install -g @marp-team/marp-cli
- Chromium/Puppeteer: 用于 PDF 导出
"""
import asyncio
import logging
import uuid
from dataclasses import dataclass
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


class ExportFormat(str, Enum):
    """导出格式"""
    HTML = "html"
    PDF = "pdf"
    PPTX = "pptx"
    PNG = "png"  # 导出为图片


class MarpTheme(str, Enum):
    """Marp 内置主题"""
    DEFAULT = "default"
    GAIA = "gaia"
    UNCOVER = "uncover"


@dataclass
class PPTRenderConfig:
    """渲染配置"""
    theme: MarpTheme = MarpTheme.DEFAULT
    paginate: bool = True
    header: str | None = None
    footer: str | None = None
    background_color: str | None = None
    color: str | None = None

    # PDF 特定
    pdf_outlines: bool = True
    pdf_notes: bool = False

    # 图片特定
    images_scale: float = 2.0

    # 自定义 CSS
    custom_css: str | None = None


@dataclass
class RenderResult:
    """渲染结果"""
    success: bool
    output_path: str | None = None
    preview_url: str | None = None
    download_url: str | None = None
    file_size: int | None = None
    error: str | None = None
    render_time: float = 0.0


class PPTRendererService:
    """PPT 渲染服务

    使用 Marp CLI 进行渲染，支持：
    - HTML 预览
    - PDF 导出
    - PPTX 导出
    - 自定义主题
    """

    def __init__(
        self,
        workspace_dir: str = "/tmp/tokendance/ppt",
        static_url_prefix: str = "/static/ppt",
        marp_bin: str = "marp"
    ):
        """初始化渲染器

        Args:
            workspace_dir: 工作目录，存放生成的文件
            static_url_prefix: 静态文件 URL 前缀
            marp_bin: Marp CLI 二进制路径
        """
        self.workspace_dir = Path(workspace_dir)
        self.workspace_dir.mkdir(parents=True, exist_ok=True)
        self.static_url_prefix = static_url_prefix
        self.marp_bin = marp_bin

        # 主题目录
        self.themes_dir = self.workspace_dir / "themes"
        self.themes_dir.mkdir(exist_ok=True)

        # 输出目录
        self.output_dir = self.workspace_dir / "output"
        self.output_dir.mkdir(exist_ok=True)

        # 初始化自定义主题
        self._init_custom_themes()

    def _init_custom_themes(self) -> None:
        """初始化自定义主题"""
        # Business 主题
        business_theme = """/* @theme business */
@import 'default';

section {
    background: linear-gradient(135deg, #1a1a2e 0%, #16213e 100%);
    color: #ffffff;
}

h1, h2 {
    color: #e94560;
}

h1 {
    font-size: 2.5em;
    font-weight: 700;
}

h2 {
    font-size: 1.8em;
    border-bottom: 3px solid #e94560;
    padding-bottom: 0.3em;
}

a {
    color: #4a90e2;
}

code {
    background: rgba(255, 255, 255, 0.1);
    color: #e94560;
}

blockquote {
    border-left: 4px solid #e94560;
    background: rgba(233, 69, 96, 0.1);
    padding: 1em;
}

section.lead {
    display: flex;
    flex-direction: column;
    justify-content: center;
    align-items: center;
    text-align: center;
}

section.lead h1 {
    font-size: 3em;
}
"""

        # Tech 主题
        tech_theme = """/* @theme tech */
@import 'default';

section {
    background: #0a0a0a;
    color: #00ff88;
    font-family: 'JetBrains Mono', 'Fira Code', monospace;
}

h1, h2 {
    color: #00ff88;
    text-shadow: 0 0 10px rgba(0, 255, 136, 0.5);
}

h1 {
    font-size: 2.5em;
}

code {
    background: rgba(0, 255, 136, 0.1);
    color: #ffffff;
    border: 1px solid #00ff88;
}

pre {
    background: #1a1a1a;
    border: 1px solid #00ff88;
}

a {
    color: #00ccff;
}

section.lead {
    background: radial-gradient(circle at center, #1a1a2e 0%, #0a0a0a 100%);
}

section.lead h1 {
    font-size: 3.5em;
    animation: glow 2s ease-in-out infinite alternate;
}

@keyframes glow {
    from { text-shadow: 0 0 10px rgba(0, 255, 136, 0.5); }
    to { text-shadow: 0 0 20px rgba(0, 255, 136, 0.8), 0 0 30px rgba(0, 255, 136, 0.6); }
}
"""

        # Minimal 主题
        minimal_theme = """/* @theme minimal */
@import 'default';

section {
    background: #fafafa;
    color: #1a1a1a;
    font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
}

h1, h2 {
    color: #111827;
    font-weight: 600;
}

h1 {
    font-size: 2.2em;
}

h2 {
    font-size: 1.5em;
}

a {
    color: #2563eb;
    text-decoration: none;
}

code {
    background: #f1f5f9;
    color: #475569;
    padding: 0.2em 0.4em;
    border-radius: 4px;
}

blockquote {
    border-left: 3px solid #e5e7eb;
    color: #6b7280;
    padding-left: 1em;
}

section.lead {
    background: #ffffff;
}

section.lead h1 {
    font-size: 3em;
    color: #111827;
}
"""

        # 保存主题文件
        themes = {
            "business.css": business_theme,
            "tech.css": tech_theme,
            "minimal.css": minimal_theme
        }

        for filename, content in themes.items():
            theme_path = self.themes_dir / filename
            if not theme_path.exists():
                theme_path.write_text(content, encoding="utf-8")
                logger.info(f"Created custom theme: {filename}")

    async def check_marp_available(self) -> bool:
        """检查 Marp CLI 是否可用"""
        try:
            proc = await asyncio.create_subprocess_exec(
                self.marp_bin, "--version",
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, _ = await proc.communicate()
            if proc.returncode == 0:
                version = stdout.decode().strip()
                logger.info(f"Marp CLI available: {version}")
                return True
        except Exception as e:
            logger.warning(f"Marp CLI not available: {e}")
        return False

    async def render_to_html(
        self,
        markdown_content: str,
        config: PPTRenderConfig | None = None,
        output_name: str | None = None
    ) -> RenderResult:
        """渲染 Markdown 为 HTML

        Args:
            markdown_content: Marp Markdown 内容
            config: 渲染配置
            output_name: 输出文件名（不含扩展名）

        Returns:
            RenderResult: 渲染结果
        """
        start_time = datetime.now()
        config = config or PPTRenderConfig()
        output_name = output_name or str(uuid.uuid4())[:8]

        # 创建临时 Markdown 文件
        md_path = self.output_dir / f"{output_name}.md"
        html_path = self.output_dir / f"{output_name}.html"

        try:
            # 写入 Markdown
            md_path.write_text(markdown_content, encoding="utf-8")

            # 构建 Marp 命令
            cmd = [
                self.marp_bin,
                str(md_path),
                "-o", str(html_path),
                "--html"
            ]

            # 添加主题
            if config.theme != MarpTheme.DEFAULT:
                theme_file = self.themes_dir / f"{config.theme.value}.css"
                if theme_file.exists():
                    cmd.extend(["--theme", str(theme_file)])

            # 添加自定义 CSS
            if config.custom_css:
                css_path = self.output_dir / f"{output_name}_custom.css"
                css_path.write_text(config.custom_css, encoding="utf-8")
                cmd.extend(["--theme", str(css_path)])

            # 执行渲染
            proc = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await proc.communicate()

            if proc.returncode != 0:
                error_msg = stderr.decode() if stderr else "Unknown error"
                logger.error(f"Marp render failed: {error_msg}")
                return RenderResult(
                    success=False,
                    error=error_msg
                )

            # 计算文件大小
            file_size = html_path.stat().st_size if html_path.exists() else 0

            render_time = (datetime.now() - start_time).total_seconds()

            return RenderResult(
                success=True,
                output_path=str(html_path),
                preview_url=f"{self.static_url_prefix}/{output_name}.html",
                file_size=file_size,
                render_time=render_time
            )

        except Exception as e:
            logger.error(f"Render error: {e}")
            return RenderResult(success=False, error=str(e))

    async def export_to_pdf(
        self,
        markdown_content: str,
        config: PPTRenderConfig | None = None,
        output_name: str | None = None
    ) -> RenderResult:
        """导出为 PDF

        Args:
            markdown_content: Marp Markdown 内容
            config: 渲染配置
            output_name: 输出文件名

        Returns:
            RenderResult: 导出结果
        """
        start_time = datetime.now()
        config = config or PPTRenderConfig()
        output_name = output_name or str(uuid.uuid4())[:8]

        md_path = self.output_dir / f"{output_name}.md"
        pdf_path = self.output_dir / f"{output_name}.pdf"

        try:
            md_path.write_text(markdown_content, encoding="utf-8")

            cmd = [
                self.marp_bin,
                str(md_path),
                "-o", str(pdf_path),
                "--pdf",
                "--allow-local-files"
            ]

            # PDF 特定选项
            if config.pdf_outlines:
                cmd.append("--pdf-outlines")
            if config.pdf_notes:
                cmd.append("--pdf-notes")

            # 主题
            if config.theme != MarpTheme.DEFAULT:
                theme_file = self.themes_dir / f"{config.theme.value}.css"
                if theme_file.exists():
                    cmd.extend(["--theme", str(theme_file)])

            proc = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await proc.communicate()

            if proc.returncode != 0:
                error_msg = stderr.decode() if stderr else "Unknown error"
                logger.error(f"PDF export failed: {error_msg}")
                return RenderResult(success=False, error=error_msg)

            file_size = pdf_path.stat().st_size if pdf_path.exists() else 0
            render_time = (datetime.now() - start_time).total_seconds()

            return RenderResult(
                success=True,
                output_path=str(pdf_path),
                download_url=f"{self.static_url_prefix}/{output_name}.pdf",
                file_size=file_size,
                render_time=render_time
            )

        except Exception as e:
            logger.error(f"PDF export error: {e}")
            return RenderResult(success=False, error=str(e))

    async def export_to_pptx(
        self,
        markdown_content: str,
        config: PPTRenderConfig | None = None,
        output_name: str | None = None
    ) -> RenderResult:
        """导出为 PPTX

        注意：Marp 原生不支持 PPTX，这里通过 HTML → PPTX 转换实现
        需要安装额外依赖：python-pptx

        Args:
            markdown_content: Marp Markdown 内容
            config: 渲染配置
            output_name: 输出文件名

        Returns:
            RenderResult: 导出结果
        """
        start_time = datetime.now()
        config = config or PPTRenderConfig()
        output_name = output_name or str(uuid.uuid4())[:8]

        try:
            # 尝试使用 python-pptx
            from pptx import Presentation
            from pptx.util import Inches, Pt

            # 解析 Markdown 为幻灯片
            slides_data = self._parse_markdown_slides(markdown_content)

            # 创建 PPTX
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9
            prs.slide_height = Inches(7.5)

            for slide_data in slides_data:
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)

                # 添加标题
                if slide_data.get("title"):
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5),
                        Inches(12), Inches(1.5)
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = slide_data["title"]
                    title_frame.paragraphs[0].font.size = Pt(44)
                    title_frame.paragraphs[0].font.bold = True

                # 添加内容
                if slide_data.get("content"):
                    content_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(2),
                        Inches(12), Inches(5)
                    )
                    content_frame = content_box.text_frame
                    content_frame.text = slide_data["content"]
                    content_frame.paragraphs[0].font.size = Pt(24)

            # 保存
            pptx_path = self.output_dir / f"{output_name}.pptx"
            prs.save(str(pptx_path))

            file_size = pptx_path.stat().st_size
            render_time = (datetime.now() - start_time).total_seconds()

            return RenderResult(
                success=True,
                output_path=str(pptx_path),
                download_url=f"{self.static_url_prefix}/{output_name}.pptx",
                file_size=file_size,
                render_time=render_time
            )

        except ImportError:
            logger.warning("python-pptx not installed, PPTX export not available")
            return RenderResult(
                success=False,
                error="PPTX export requires python-pptx package. Install with: pip install python-pptx"
            )
        except Exception as e:
            logger.error(f"PPTX export error: {e}")
            return RenderResult(success=False, error=str(e))

    def _parse_markdown_slides(self, markdown_content: str) -> list[dict[str, Any]]:
        """解析 Markdown 为幻灯片数据"""
        slides = []

        # 移除 frontmatter
        content = markdown_content
        if content.startswith("---"):
            parts = content.split("---", 2)
            if len(parts) >= 3:
                content = parts[2]

        # 按 --- 分割幻灯片
        slide_texts = content.split("\n---\n")

        for slide_text in slide_texts:
            slide_text = slide_text.strip()
            if not slide_text:
                continue

            slide_data = {"title": "", "content": ""}
            lines = slide_text.split("\n")

            for line in lines:
                line = line.strip()
                # 跳过 Marp 指令
                if line.startswith("<!--") or line.startswith("```"):
                    continue

                # 提取标题
                if line.startswith("# "):
                    slide_data["title"] = line[2:].strip()
                elif line.startswith("## "):
                    if not slide_data["title"]:
                        slide_data["title"] = line[3:].strip()
                    else:
                        slide_data["content"] += line[3:] + "\n"
                elif line:
                    slide_data["content"] += line + "\n"

            slide_data["content"] = slide_data["content"].strip()
            slides.append(slide_data)

        return slides

    async def export(
        self,
        markdown_content: str,
        format: ExportFormat,
        config: PPTRenderConfig | None = None,
        output_name: str | None = None
    ) -> RenderResult:
        """统一导出接口

        Args:
            markdown_content: Marp Markdown 内容
            format: 导出格式
            config: 渲染配置
            output_name: 输出文件名

        Returns:
            RenderResult: 导出结果
        """
        if format == ExportFormat.HTML:
            return await self.render_to_html(markdown_content, config, output_name)
        elif format == ExportFormat.PDF:
            return await self.export_to_pdf(markdown_content, config, output_name)
        elif format == ExportFormat.PPTX:
            return await self.export_to_pptx(markdown_content, config, output_name)
        else:
            return RenderResult(success=False, error=f"Unsupported format: {format}")

    def cleanup_old_files(self, max_age_hours: int = 24) -> int:
        """清理旧文件

        Args:
            max_age_hours: 最大保留时间（小时）

        Returns:
            int: 删除的文件数
        """
        import time

        deleted = 0
        cutoff = time.time() - (max_age_hours * 3600)

        for path in self.output_dir.iterdir():
            if path.is_file() and path.stat().st_mtime < cutoff:
                try:
                    path.unlink()
                    deleted += 1
                except Exception as e:
                    logger.warning(f"Failed to delete {path}: {e}")

        logger.info(f"Cleaned up {deleted} old PPT files")
        return deleted


# ==================== 单例实例 ====================

_renderer_instance: PPTRendererService | None = None


def get_ppt_renderer() -> PPTRendererService:
    """获取 PPT 渲染器单例"""
    global _renderer_instance
    if _renderer_instance is None:
        _renderer_instance = PPTRendererService()
    return _renderer_instance


async def check_renderer_health() -> dict[str, Any]:
    """检查渲染器健康状态"""
    renderer = get_ppt_renderer()
    marp_available = await renderer.check_marp_available()

    return {
        "service": "ppt_renderer",
        "status": "healthy" if marp_available else "degraded",
        "marp_cli": marp_available,
        "workspace": str(renderer.workspace_dir),
        "themes_available": list(renderer.themes_dir.glob("*.css"))
    }
